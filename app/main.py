import os
import uuid
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from pydantic import BaseModel,EmailStr
from vector_store import add_document, search, embed_text, get_user_collection
from cerebras.cloud.sdk import Cerebras
from PyPDF2 import PdfReader
import mysql.connector
import random
import string


import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

from docx import Document
from pptx import Presentation
import pandas as pd
import pytesseract
from PIL import Image

app = FastAPI(title="Cerebras RAG AI Assistant")
client = Cerebras(api_key=os.environ.get("CEREBRAS_API_KEY"))

class ChatRequest(BaseModel):
    user_id: str
    query: str

# Utility: split text into chunks
def chunk_text(text, chunk_size=500):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]


# ----------------------------
# Upload file (user-specific)
# ----------------------------
# @app.post("/upload-file")
# async def upload_file(
#     user_id: str = Form(...), 
#     file: UploadFile = File(...)
# ):
#     """
#     Upload a file for a specific user, extract text, chunk it, and store embeddings.
#     """
#     # Ensure storage directory exists
#     storage_dir = f"app/file_storage/{user_id}"
#     os.makedirs(storage_dir, exist_ok=True)

#     # Save file
#     file_id = str(uuid.uuid4())
#     file_path = os.path.join(storage_dir, f"{file_id}_{file.filename}")
#     with open(file_path, "wb") as f:
#         f.write(await file.read())

#     # Extract text
#     text = ""
#     if file.filename.endswith(".pdf"):
#         reader = PdfReader(file_path)
#         text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
#     else:
#         with open(file_path, "r", errors="ignore") as f:
#             text = f.read()

#     # Chunk + store (user-specific)
#     chunks = chunk_text(text)
#     add_document(user_id, file_id, chunks, source=file.filename)

#     return {
#         "message": "File uploaded and indexed",
#         "file_id": file_id,
#         "source": file.filename,
#         "user_id": user_id
#     }

@app.post("/upload-file")
async def upload_file(
    user_id: str = Form(...), 
    file: UploadFile = File(...)
):
    """
    Upload a file for a specific user, extract text, chunk it, and store embeddings.
    Supports PDF, TXT, DOCX, PPTX, CSV, XLSX, and images (OCR).
    """
    # Ensure storage directory exists
    storage_dir = f"app/file_storage/{user_id}"
    os.makedirs(storage_dir, exist_ok=True)

    # Save file
    file_id = str(uuid.uuid4())
    file_path = os.path.join(storage_dir, f"{file_id}_{file.filename}")
    with open(file_path, "wb") as f:
        f.write(await file.read())

    # Extract text
    text = ""
    filename_lower = file.filename.lower()
    
    try:
        if filename_lower.endswith(".pdf"):
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
        
        elif filename_lower.endswith(".txt"):
            with open(file_path, "r", errors="ignore") as f:
                text = f.read()

        elif filename_lower.endswith(".docx"):
            doc = Document(file_path)
            text = " ".join([para.text for para in doc.paragraphs])

        elif filename_lower.endswith(".pptx"):
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            text = " ".join(slides_text)

        elif filename_lower.endswith((".csv", ".xlsx")):
            if filename_lower.endswith(".csv"):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            text = df.astype(str).apply(lambda x: " ".join(x), axis=1).str.cat(sep=" ")

        elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)

        else:
            return {"error": "Unsupported file type"}
    
    except Exception as e:
        return {"error": f"Failed to extract text: {str(e)}"}

    # Chunk + store (user-specific)
    chunks = chunk_text(text)
    add_document(user_id, file_id, chunks, source=file.filename)

    return {
        "message": "File uploaded and indexed",
        "file_id": file_id,
        "source": file.filename,
        "user_id": user_id
    }

# ----------------------------
# Chat endpoint (user-specific)
# ----------------------------
# @app.post("/chat")
# def chat(request: ChatRequest):
#     """
#     Search documents for the user and respond using Cerebras LLM.
#     """
#     # Search vector DB scoped to user
#     results = search(request.user_id, request.query, top_k=3)
#     retrieved_docs = results["documents"][0]
#     sources = results["metadatas"][0]

#     context_text = "\n".join(
#         [f"Source: {src['source']}\nContent: {doc}" for doc, src in zip(retrieved_docs, sources)]
#     )

#     # System prompt
#     SYSTEM_PROMPT = """
# You are a helpful and knowledgeable assistant. Your responses must follow these rules STRICTLY:

# **WHEN THE ANSWER IS IN THE CONTEXT:**
# - Provide a clear, specific answer using ONLY the information from the provided context
# - Always cite your source by mentioning which document the information came from
# - Write in a friendly, casual tone like a knowledgeable person explaining something
# - Be concise but thorough - give the complete answer found in the context

# **WHEN THE ANSWER IS NOT IN THE CONTEXT:**
# - DO NOT try to make up an answer or use outside knowledge
# - Politely state that you couldn't find the specific information in the provided documents
# - Offer a friendly suggestion for how the user might find the information
# - Keep it warm and human-like - don't sound robotic or apologetic

# Remember: Your knowledge is limited to exactly what's in the provided context documents.
# """

#     # Cerebras chat call
#     response = client.chat.completions.create(
#         model="llama-4-scout-17b-16e-instruct",
#         messages=[
#             {"role": "system", "content": SYSTEM_PROMPT},
#             {"role": "user", "content": f"""Here is my question: {request.query}

# Here are the documents I have for context:
# {context_text}

# Please answer my question using ONLY the information above. If the answer isn't there, just let me know politely."""}
#         ]
#     )

#     return {
#         "answer": response.choices[0].message.content,
#         "sources": [src['source'] for src in sources],
#         "user_id": request.user_id
#     }


@app.post("/chat")
def chat(request: ChatRequest):
    """
    Search documents for the user and respond using Cerebras LLM.
    """
    # Check if the query is a greeting or casual conversation
    greeting_keywords = ["hello", "hi", "hey", "good morning", "good afternoon", "good evening", 
                        "how are you", "what's up", "hey there", "hi there"]
    
    is_greeting = any(keyword in request.query.lower() for keyword in greeting_keywords)
    
    if is_greeting:
        return handle_greeting(request.query, request.user_id)
    
    # Search vector DB scoped to user
    results = search(request.user_id, request.query, top_k=3)
    retrieved_docs = results["documents"][0]
    sources = results["metadatas"][0]

    context_text = "\n".join(
        [f"Source: {src['source']}\nContent: {doc}" for doc, src in zip(retrieved_docs, sources)]
    )

    # System prompt
    SYSTEM_PROMPT = """
You are a helpful and knowledgeable assistant. Your responses must follow these rules STRICTLY:

**WHEN THE ANSWER IS IN THE CONTEXT:**
- Provide a clear, specific answer using ONLY the information from the provided context
- Always cite your source by mentioning which document the information came from
- Write in a friendly, casual tone like a knowledgeable person explaining something
- Be concise but thorough - give the complete answer found in the context

**WHEN THE ANSWER IS NOT IN THE CONTEXT:**
- DO NOT try to make up an answer or use outside knowledge
- Politely state that you couldn't find the specific information in the provided documents
- Offer a friendly suggestion for how the user might find the information
- Keep it warm and human-like - don't sound robotic or apologetic

**FOR GREETINGS AND CASUAL CONVERSATION:**
- Be warm, friendly, and engaging
- Keep it brief but personable
- Naturally transition to offering help with their questions

Remember: Your knowledge is limited to exactly what's in the provided context documents.
"""

    # Cerebras chat call
    response = client.chat.completions.create(
        model="llama-4-scout-17b-16e-instruct",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"""Here is my question: {request.query}

Here are the documents I have for context:
{context_text}

Please answer my question using ONLY the information above. If the answer isn't there, just let me know politely."""}
        ]
    )

    return {
        "answer": response.choices[0].message.content,
        "sources": [src['source'] for src in sources],
        "user_id": request.user_id
    }


def handle_greeting(query: str, user_id: str):
    """Handle greetings and casual conversation without document search"""
    
    greeting_responses = [
        "Hi there! 👋 I'm here to help you find information from your documents. What can I help you with today?",
        "Hello! 😊 I'm ready to search through your documents and answer your questions. What would you like to know?",
        "Hey! Great to see you. I'm here to help you find information - just ask me anything about your documents!",
        "Hi! I'm your document assistant. I can search through your files and answer questions based on them. What would you like to know?"
    ]
    
    # More specific responses for certain greetings
    query_lower = query.lower()
    if "how are you" in query_lower:
        response = "I'm doing great, thanks for asking! 😊 Ready to help you find information from your documents. What can I help you with?"
    elif "what's up" in query_lower or "sup" in query_lower:
        response = "Not much! Just here and ready to search your documents for you. What information are you looking for?"
    elif "good morning" in query_lower:
        response = "Good morning! ☀️ Hope you're having a great start to your day. What can I help you find in your documents?"
    elif "good afternoon" in query_lower:
        response = "Good afternoon! 😊 Ready to help you find whatever information you need from your documents."
    elif "good evening" in query_lower:
        response = "Good evening! 🌙 I'm here to help you search through your documents. What would you like to know?"
    else:
        # Random selection from general greetings
        import random
        response = random.choice(greeting_responses)
    
    return {
        "answer": response,
        "sources": [],
        "user_id": user_id
    }

# ----------------------------
# Database Connection
# ----------------------------
def get_db():
    return mysql.connector.connect(
        host=os.getenv("MYSQL_HOST", "localhost"),
        user=os.getenv("MYSQL_USER", "root"),
        password=os.getenv("MYSQL_PASSWORD", "root"),
        database=os.getenv("MYSQL_DB", "authdb")
    )

# ----------------------------
# Auth Models
# ----------------------------
class OTPRequest(BaseModel):
    email: EmailStr

class VerifyRequest(BaseModel):
    user_id: str
    otp: str

# ----------------------------
# OTP Helpers
# ----------------------------
def generate_otp(length=6):
    return ''.join(random.choices(string.digits, k=length))

def send_email(to_email: str, otp: str):
    sender_email = os.getenv("EMAIL_USER", "")
    sender_password =  os.getenv("EMAIL_PASS", "") 

    subject = "Your OTP Code"
    body = f"Your One-Time Password (OTP) is: {otp}\n\nIt will expire in 10 minutes."

    msg = MIMEMultipart()
    msg["From"] = sender_email
    msg["To"] = to_email
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain"))

    try:
        # Gmail SMTP setup
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()  # Secure the connection
            server.login(sender_email, sender_password)
            server.send_message(msg)
            print(f"✅ OTP sent successfully to {to_email}")
    except Exception as e:
        print(f"❌ Error sending email: {e}")


@app.post("/send-otp")
def send_otp(req: OTPRequest):
    conn = get_db()
    cursor = conn.cursor()

    otp = generate_otp()

    # Check if email already exists
    cursor.execute("SELECT id FROM otps WHERE email = %s", (req.email,))
    row = cursor.fetchone()   # ✅ consume result here

    if row:
        user_id = row[0]
        # Clear any pending results just in case
        cursor.fetchall()  # ✅ ensures no unread results
        cursor.execute("UPDATE otps SET otp = %s WHERE id = %s", (otp, user_id))
        message = "OTP updated for existing user"
    else:
        cursor.execute("INSERT INTO otps (email, otp) VALUES (%s, %s)", (req.email, otp))
        user_id = cursor.lastrowid
        message = "OTP generated for new user"

    conn.commit()
    cursor.close()
    conn.close()

    send_email(req.email, otp)

    return {
        "status": "success",
        "user_id": str(user_id),
        "email": req.email,
        "otp": otp,
        "message": message
    }


@app.post("/verify-otp")
def verify_otp(req: VerifyRequest):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT otp FROM otps WHERE id=%s ORDER BY created_at DESC LIMIT 1", (req.user_id,))
    row = cursor.fetchone()

    cursor.close()
    conn.close()

    if not row:
        raise HTTPException(status_code=404, detail="OTP not found")

    if row[0] == req.otp:
        return {"status": "success", "message": "OTP verified"}
    else:
        raise HTTPException(status_code=400, detail="Invalid OTP")
